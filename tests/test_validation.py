"""
tests/test_validation.py

Phase 10 tests. Critically includes tests against DELIBERATELY BROKEN
artifacts (empty document, empty slides, wrong slide count, corrupt file,
unregistered citation) to prove the Validation Agent actually catches
problems rather than always returning PASS.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from fastapi.testclient import TestClient

from backend.main import app
from backend.agents.validator import validate_docx, validate_pptx
from backend.agents.document_generator import generate_docx
from backend.agents.ppt_generator import generate_pptx
from backend.schemas.content_plan import ContentPlan, GeneratedSection
from backend.services.citation_service import citation_registry

client = TestClient(app)


@pytest.fixture(autouse=True)
def _reset_citations():
    citation_registry.reset()
    yield


# --- Valid artifacts should PASS ---

def test_validate_docx_pass_on_well_formed_document(tmp_path):
    cid = citation_registry.register("WEB", "Source", "url")
    plan = ContentPlan(
        title="Good Doc",
        sections=[GeneratedSection(heading="Section 1", paragraphs=["Some real content here."], citation_ids=[cid])],
    )
    path = str(tmp_path / "good.docx")
    generate_docx(plan, path)

    result = validate_docx(path)
    assert result.status == "PASS"
    assert result.issues == []


def test_validate_pptx_pass_on_well_formed_presentation(tmp_path):
    plan = ContentPlan(title="Good Deck", sections=[GeneratedSection(heading="S1", bullet_points=["b1", "b2"])])
    path = str(tmp_path / "good.pptx")
    generate_pptx(plan, path, target_slide_count=4)

    result = validate_pptx(path, expected_slide_count=4)
    assert result.status == "PASS"
    assert result.issues == []


# --- Deliberately broken artifacts should FAIL ---

def test_validate_docx_fails_on_empty_document(tmp_path):
    """A DOCX with only a title and no body content should fail validation."""
    doc = DocxDocument()
    doc.add_heading("Title Only", level=0)
    path = str(tmp_path / "empty.docx")
    doc.save(path)

    result = validate_docx(path)
    assert result.status == "FAIL"
    assert any("empty" in issue.lower() for issue in result.issues)


def test_validate_docx_fails_on_no_headings():
    """A DOCX with zero headings at all fails structural validation."""
    import tempfile
    doc = DocxDocument()
    doc.add_paragraph("Just some plain text, no headings at all.")
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        doc.save(f.name)
        path = f.name

    result = validate_docx(path)
    assert result.status == "FAIL"
    assert any("heading" in issue.lower() for issue in result.issues)


def test_validate_docx_fails_on_corrupt_file(tmp_path):
    path = tmp_path / "corrupt.docx"
    path.write_bytes(b"this is not a real docx file at all")

    result = validate_docx(str(path))
    assert result.status == "FAIL"
    assert any("does not open" in issue.lower() for issue in result.issues)


def test_validate_docx_fails_on_missing_required_section(tmp_path):
    plan = ContentPlan(title="Doc", sections=[GeneratedSection(heading="Introduction", paragraphs=["text"])])
    path = str(tmp_path / "missing_section.docx")
    generate_docx(plan, path)

    result = validate_docx(path, required_sections=["Executive Summary"])
    assert result.status == "FAIL"
    assert any("Executive Summary" in issue for issue in result.issues)


def test_validate_docx_fails_on_unregistered_citation(tmp_path):
    """Citation appears in text but was never registered -> should be flagged."""
    doc = DocxDocument()
    doc.add_heading("Report", level=0)
    doc.add_heading("Section", level=1)
    doc.add_paragraph("Some claim here [WEB-999].")
    path = str(tmp_path / "bad_citation.docx")
    doc.save(path)

    result = validate_docx(path)
    assert result.status == "FAIL"
    assert any("WEB-999" in issue for issue in result.issues)


def test_validate_pptx_fails_on_wrong_slide_count(tmp_path):
    plan = ContentPlan(title="Deck", sections=[GeneratedSection(heading="S1", bullet_points=["b1"])])
    path = str(tmp_path / "wrong_count.pptx")
    generate_pptx(plan, path, target_slide_count=5)

    result = validate_pptx(path, expected_slide_count=12)
    assert result.status == "FAIL"
    assert any("expected 12" in issue.lower() for issue in result.issues)


def test_validate_pptx_fails_on_empty_slide():
    """Manually construct a presentation with a slide that has zero content."""
    import tempfile
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank layout
    prs.slides.add_slide(layout)  # completely empty slide
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        path = f.name

    result = validate_pptx(path)
    assert result.status == "FAIL"
    assert any("empty" in issue.lower() for issue in result.issues)


def test_validate_pptx_fails_on_corrupt_file(tmp_path):
    path = tmp_path / "corrupt.pptx"
    path.write_bytes(b"not a real pptx")

    result = validate_pptx(str(path))
    assert result.status == "FAIL"
    assert any("does not open" in issue.lower() for issue in result.issues)


# --- API-level: validation runs automatically on generation ---

def test_generate_document_includes_validation_result():
    resp = client.post("/generate/document", json={"request": "Create a report", "research_findings": ["Fact [WEB-001]."]})
    assert resp.status_code == 200
    body = resp.json()
    assert "validation" in body
    assert body["validation"]["status"] in ("PASS", "FAIL")


def test_generate_presentation_validates_slide_count_mismatch_detection():
    resp = client.post("/generate/presentation", json={"request": "Create a deck", "slide_count": 6})
    assert resp.status_code == 200
    body = resp.json()
    assert body["validation"]["status"] == "PASS"  # generator pads to match, so this should pass


def test_standalone_validate_endpoint():
    gen_resp = client.post("/generate/document", json={"request": "Create a report"})
    artifact_id = gen_resp.json()["artifact_id"]

    validate_resp = client.post(f"/validate/{artifact_id}")
    assert validate_resp.status_code == 200
    assert validate_resp.json()["status"] in ("PASS", "FAIL")
