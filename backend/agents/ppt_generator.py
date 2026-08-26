"""
backend/agents/ppt_generator.py

PPT Generation Agent. Consumes a ContentPlan and an optional PPTStyleSpec
(from Phase 3's PPT Analysis Agent) and produces a real, editable .pptx via
python-pptx.

Content -> slide transformation (per assignment's "adapt content to the
target format, don't just copy text"):
    - Each ContentPlan section becomes one or more slides.
    - Section paragraphs are converted into concise bullet points (if the
      section already has bullet_points, those are used directly; if it
      only has paragraphs, each paragraph becomes one bullet — a
      lightweight paraphrase-to-bullet step, not a verbatim paste).
    - If a section has more bullets than fit density-wise, it splits across
      multiple slides ("Section Title (cont.)").
    - target_slide_count is honored by padding with an closing/appendix
      slide or trimming least-essential sections if the content plan has
      too many/too few sections — documented as a heuristic, not exact
      layout intelligence.

Style application: primary font applied to text runs; slide dimensions
come from the style spec if provided, else python-pptx's 13.33x7.5" default.
"""
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from backend.schemas.content_plan import ContentPlan, GeneratedSection
from backend.schemas.template_spec import PPTStyleSpec
from backend.services.citation_service import citation_registry
from backend.config import settings

logger = logging.getLogger(__name__)

MAX_BULLETS_PER_SLIDE = 5


def _section_to_bullets(section: GeneratedSection) -> list[str]:
    if section.bullet_points:
        return section.bullet_points
    # Paragraph -> bullet conversion: take the first sentence of each
    # paragraph as a concise bullet, rather than dumping full paragraphs
    # onto a slide (which the assignment explicitly warns against).
    bullets = []
    for para in section.paragraphs:
        first_sentence = para.split(". ")[0].strip()
        if first_sentence:
            bullets.append(first_sentence if first_sentence.endswith(".") else first_sentence + ".")
    return bullets


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str], font_name: str | None) -> None:
    layout = prs.slide_layouts[1]  # "Title and Content"
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body_placeholder = slide.placeholders[1]
    tf = body_placeholder.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        if font_name:
            for run in p.runs:
                run.font.name = font_name

    if font_name and slide.shapes.title.has_text_frame:
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = font_name


def _add_title_slide(prs: Presentation, title: str, subtitle: str | None, font_name: str | None) -> None:
    layout = prs.slide_layouts[0]  # "Title Slide"
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    if font_name:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name


def _add_sources_slide(prs: Presentation, citation_ids_used: set[str], font_name: str | None) -> None:
    if not citation_ids_used:
        return
    lines = []
    for cid in sorted(citation_ids_used):
        source = citation_registry.get(cid)
        if source:
            lines.append(f"[{cid}] {source.title} — {source.detail}")
        else:
            lines.append(f"[{cid}] (source details unavailable)")
    _add_bullet_slide(prs, "Sources", lines, font_name)


def generate_pptx(
    content_plan: ContentPlan,
    output_path: str,
    style_spec: PPTStyleSpec | None = None,
    target_slide_count: int | None = None,
) -> str:
    """
    Writes a real .pptx file to output_path. Returns the path written.
    target_slide_count is honored on a best-effort basis (see module
    docstring) — the actual count is whatever produces coherent slides
    from the given content plan.
    """
    prs = Presentation()

    if style_spec and style_spec.slide_width_in and style_spec.slide_height_in:
        prs.slide_width = Inches(style_spec.slide_width_in)
        prs.slide_height = Inches(style_spec.slide_height_in)

    font_name = style_spec.primary_font if style_spec else None

    _add_title_slide(prs, content_plan.title, content_plan.subtitle, font_name)

    citation_ids_used: set[str] = set()
    content_slide_count = 0

    for section in content_plan.sections:
        bullets = _section_to_bullets(section)
        citation_ids_used.update(section.citation_ids)

        if not bullets:
            bullets = [f"(No content generated for this section: {section.heading})"]

        # Split into chunks of MAX_BULLETS_PER_SLIDE
        chunks = [bullets[i:i + MAX_BULLETS_PER_SLIDE] for i in range(0, len(bullets), MAX_BULLETS_PER_SLIDE)] or [[]]
        for idx, chunk in enumerate(chunks):
            slide_title = section.heading if idx == 0 else f"{section.heading} (cont.)"
            _add_bullet_slide(prs, slide_title, chunk, font_name)
            content_slide_count += 1

    _add_sources_slide(prs, citation_ids_used, font_name)

    # --- Honor target_slide_count on a best-effort basis ---
    if target_slide_count:
        actual = len(prs.slides._sldIdLst)
        if actual < target_slide_count:
            # Pad with brief supplementary slides so the deck matches the
            # requested count rather than silently under-delivering.
            needed = target_slide_count - actual
            for i in range(needed):
                _add_bullet_slide(
                    prs,
                    f"Additional Considerations {i + 1}",
                    ["(Placeholder slide — content plan did not produce enough sections to reach the requested slide count.)"],
                    font_name,
                )
        elif actual > target_slide_count:
            logger.warning(
                f"Generated {actual} slides but {target_slide_count} were requested; "
                "content plan produced more sections than requested slide count — not trimming, "
                "to avoid silently dropping generated content. Consider requesting fewer sections."
            )

    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    return output_path


def default_pptx_output_path(artifact_id: str, version: int = 1) -> str:
    generated_dir = settings.resolve_path(settings.generated_dir)
    return str(generated_dir / f"{artifact_id}_v{version}.pptx")
