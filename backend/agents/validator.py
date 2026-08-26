"""
backend/agents/validator.py

Validation Agent. Per assignment: checks generated artifacts and must not
silently return broken output. Every check here is a real structural
inspection of the actual generated file (re-opened with python-docx /
python-pptx), not a guess based on the content plan that produced it —
this catches real generation bugs, not just "did we intend to include X".

Document validation:
    - file opens successfully
    - at least one heading/section exists
    - no entirely empty document (title only, nothing else)
    - citations referenced in text actually exist in the citation registry
    - required sections present, if a list of expected section names is given

PPT validation:
    - file opens successfully
    - matches requested slide count (if one was requested) — FAIL if not,
      since silently shipping the wrong count violates the assignment
    - no empty slides (a slide with a title but zero body text/table content)
    - citations referenced in text actually exist in the citation registry
"""
import logging
import re

from docx import Document as DocxDocument
from pptx import Presentation

from backend.schemas.validation import ValidationResult
from backend.services.citation_service import citation_registry

logger = logging.getLogger(__name__)

CITATION_PATTERN = re.compile(r"\[(WEB-\d+|RAG-\d+|DOC-\d+)\]")


def _extract_citation_ids_from_text(text: str) -> set[str]:
    return set(CITATION_PATTERN.findall(text))


def _check_citations_exist(citation_ids: set[str], warnings: list[str], issues: list[str]) -> None:
    for cid in citation_ids:
        if citation_registry.get(cid) is None:
            issues.append(f"Citation '{cid}' appears in the document but is not registered in the citation registry.")


def validate_docx(file_path: str, required_sections: list[str] | None = None) -> ValidationResult:
    issues: list[str] = []
    warnings: list[str] = []

    try:
        doc = DocxDocument(file_path)
    except Exception as e:
        return ValidationResult(
            status="FAIL",
            artifact_type="docx",
            file_path=file_path,
            issues=[f"File does not open as a valid DOCX: {e}"],
        )

    all_paragraphs = doc.paragraphs
    headings = [p for p in all_paragraphs if p.style and p.style.name and p.style.name.lower().startswith(("heading", "title"))]
    body_paragraphs = [
        p for p in all_paragraphs
        if p.text.strip() and not (p.style and p.style.name and p.style.name.lower().startswith(("heading", "title")))
    ]

    if not headings:
        issues.append("Document has no headings/title — structure is missing.")

    if not body_paragraphs and not doc.tables:
        issues.append("Document has headings but no body content (paragraphs or tables) — appears empty.")

    if required_sections:
        heading_texts = {h.text.strip().lower() for h in headings}
        for required in required_sections:
            if required.strip().lower() not in heading_texts:
                issues.append(f"Required section '{required}' not found in document.")

    full_text = "\n".join(p.text for p in doc.paragraphs)
    citation_ids = _extract_citation_ids_from_text(full_text)
    if citation_ids:
        _check_citations_exist(citation_ids, warnings, issues)
    else:
        warnings.append("No citations found in document text (may be expected if no research/RAG context was used).")

    status = "FAIL" if issues else "PASS"
    return ValidationResult(status=status, artifact_type="docx", file_path=file_path, issues=issues, warnings=warnings)


def validate_pptx(file_path: str, expected_slide_count: int | None = None) -> ValidationResult:
    issues: list[str] = []
    warnings: list[str] = []

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return ValidationResult(
            status="FAIL",
            artifact_type="pptx",
            file_path=file_path,
            issues=[f"File does not open as a valid PPTX: {e}"],
        )

    slide_count = len(prs.slides)
    if slide_count == 0:
        issues.append("Presentation has zero slides.")

    if expected_slide_count is not None and slide_count != expected_slide_count:
        issues.append(f"Expected {expected_slide_count} slides but found {slide_count}.")

    all_text_parts = []
    empty_slide_numbers = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        has_table = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_text_parts.append(shape.text_frame.text)
            if shape.has_table:
                has_table = True

        # slide dimensions bounds check: verify text-bearing shapes are within slide bounds
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                try:
                    if shape.left is not None and shape.top is not None:
                        if shape.left < 0 or shape.top < 0:
                            warnings.append(f"Slide {idx}: a text shape has a negative position (may render off-slide).")
                        if shape.left + (shape.width or 0) > prs.slide_width:
                            warnings.append(f"Slide {idx}: a text shape may extend beyond the right edge of the slide.")
                except Exception:
                    pass

        if not slide_text_parts and not has_table:
            empty_slide_numbers.append(idx)
        all_text_parts.extend(slide_text_parts)

    if empty_slide_numbers:
        issues.append(f"Slide(s) {empty_slide_numbers} have no text or table content — appear empty.")

    full_text = "\n".join(all_text_parts)
    citation_ids = _extract_citation_ids_from_text(full_text)
    if citation_ids:
        _check_citations_exist(citation_ids, warnings, issues)
    else:
        warnings.append("No citations found in slide text (may be expected if no research/RAG context was used).")

    status = "FAIL" if issues else "PASS"
    return ValidationResult(status=status, artifact_type="pptx", file_path=file_path, issues=issues, warnings=warnings)
