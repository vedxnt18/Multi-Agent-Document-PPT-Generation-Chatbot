"""
backend/services/ppt_service.py

Extraction for PPTX files, producing the normalized ExtractionResult schema.
Walks slides in order, separating title/body placeholders from other shapes,
tables, and speaker notes.
"""
import logging

from pptx import Presentation
from pptx.util import Emu

from backend.schemas.extraction import ExtractionResult, ContentBlock, BlockType, SourceType

logger = logging.getLogger(__name__)


def extract_pptx(file_path: str, file_id: str, original_filename: str) -> ExtractionResult:
    result = ExtractionResult(
        file_id=file_id,
        original_filename=original_filename,
        file_extension=".pptx",
        document_type="pptx",
    )

    try:
        prs = Presentation(file_path)
    except Exception as e:
        result.warnings.append(f"Failed to open PPTX: {e}")
        return result

    result.page_or_slide_count = len(prs.slides)

    for slide_index, slide in enumerate(prs.slides):
        slide_number = slide_index + 1
        title_captured = False

        for shape in slide.shapes:
            # --- Title / body text placeholders ---
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                is_title = False
                try:
                    is_title = (
                        shape.placeholder_format is not None
                        and shape.placeholder_format.type is not None
                        and "TITLE" in str(shape.placeholder_format.type)
                    )
                except Exception:
                    is_title = False

                block_type = BlockType.SLIDE_TITLE if (is_title and not title_captured) else BlockType.SLIDE_BODY
                if block_type == BlockType.SLIDE_TITLE:
                    title_captured = True

                bbox = None
                try:
                    bbox = [
                        float(Emu(shape.left).inches),
                        float(Emu(shape.top).inches),
                        float(Emu(shape.left + shape.width).inches),
                        float(Emu(shape.top + shape.height).inches),
                    ]
                except Exception:
                    bbox = None

                result.blocks.append(
                    ContentBlock(
                        block_type=block_type,
                        text=text,
                        slide_number=slide_number,
                        source_type=SourceType.NATIVE_TEXT,
                        bbox=bbox,
                        metadata={"shape_name": shape.name},
                    )
                )

            # --- Tables ---
            if shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                result.blocks.append(
                    ContentBlock(
                        block_type=BlockType.TABLE,
                        text=f"[Slide {slide_number} table with {len(rows)} rows]",
                        slide_number=slide_number,
                        source_type=SourceType.NATIVE_TEXT,
                        metadata={"rows": rows},
                    )
                )

            # --- Images (record presence, not OCR by default — slides are
            #     usually template graphics, not scanned content) ---
            if shape.shape_type is not None and str(shape.shape_type) == "PICTURE (13)":
                result.blocks.append(
                    ContentBlock(
                        block_type=BlockType.IMAGE,
                        text="",
                        slide_number=slide_number,
                        source_type=SourceType.NATIVE_TEXT,
                        metadata={"shape_name": shape.name},
                    )
                )

        # --- Speaker notes ---
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                result.blocks.append(
                    ContentBlock(
                        block_type=BlockType.SPEAKER_NOTES,
                        text=notes_text,
                        slide_number=slide_number,
                        source_type=SourceType.NATIVE_TEXT,
                    )
                )

    return result
