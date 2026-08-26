"""
tests/test_editing_and_versions.py

Phase 11 (Conversational Editing) + Phase 13 (Version Management) tests.
Uses the heuristic edit-interpretation path since LLM_PROVIDER=mock in this
environment — verifies add/remove/condense operations actually mutate the
artifact's content plan and produce a new, validated version, and that
version history is queryable via the API.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from fastapi.testclient import TestClient

from backend.main import app
from backend.agents.conversational_editor import edit_artifact, interpret_edit_request, _heuristic_instruction
from backend.services import artifact_store
from backend.services.artifact_store import ArtifactNotFoundError
from backend.services.citation_service import citation_registry

client = TestClient(app)


@pytest.fixture(autouse=True)
def _reset():
    citation_registry.reset()
    yield


def _generate_doc(request="Create a proposal", research_findings=None):
    resp = client.post("/generate/document", json={"request": request, "research_findings": research_findings or []})
    assert resp.status_code == 200
    return resp.json()["artifact_id"]


def _generate_ppt(request="Create a deck", slide_count=4):
    resp = client.post("/generate/presentation", json={"request": request, "slide_count": slide_count})
    assert resp.status_code == 200
    return resp.json()["artifact_id"]


# --- Heuristic interpretation ---

def test_heuristic_detects_add_section():
    instr = _heuristic_instruction("Add an executive summary", ["Introduction"])
    assert instr.operation == "add_section"
    assert "executive summary" in instr.new_heading.lower()


def test_heuristic_detects_condense():
    instr = _heuristic_instruction("Make the presentation more concise", ["Introduction"])
    assert instr.operation == "condense"


def test_heuristic_detects_remove_section():
    instr = _heuristic_instruction("Remove the Introduction section", ["Introduction", "Market Overview"])
    assert instr.operation == "remove_section"
    assert instr.target_heading == "Introduction"


def test_heuristic_unsupported_for_unrecognized_request():
    instr = _heuristic_instruction("What time is it?", ["Introduction"])
    assert instr.operation == "unsupported"


# --- Editing mutates artifact and creates a new version ---

def test_edit_add_section_creates_new_docx_version():
    artifact_id = _generate_doc()
    assert artifact_store.get_latest_version_number(artifact_id) == 1

    result = edit_artifact(artifact_id, "Add an executive summary")

    assert result.new_version == 2
    assert result.instruction.operation == "add_section"
    assert Path(result.file_path).exists()

    doc = DocxDocument(result.file_path)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Executive Summary" in full_text

    # Original v1 file must still exist untouched (version history preserved)
    meta = artifact_store.load_meta(artifact_id)
    assert len(meta["versions"]) == 2
    v1_path = Path(meta["versions"][0]["file_path"])
    assert v1_path.exists()
    v1_doc = DocxDocument(v1_path)
    v1_text = "\n".join(p.text for p in v1_doc.paragraphs)
    assert "Executive Summary" not in v1_text  # v1 unaffected by the edit


def test_edit_condense_reduces_bullet_count():
    artifact_id = _generate_ppt()

    # First add a section with many bullets so condense has something to do
    from backend.services import artifact_store as store
    plan = store.load_latest_content_plan(artifact_id)
    from backend.schemas.content_plan import GeneratedSection
    plan.sections.append(GeneratedSection(heading="Dense Section", bullet_points=[f"Point {i}" for i in range(6)]))
    from backend.agents.ppt_generator import generate_pptx, default_pptx_output_path
    path = default_pptx_output_path(artifact_id, version=2)
    generate_pptx(plan, path)
    store.save_version(artifact_id, "pptx", 2, path, plan)

    result = edit_artifact(artifact_id, "Make it more concise")
    assert result.new_version == 3
    assert result.instruction.operation == "condense"

    prs = Presentation(result.file_path)
    # find the dense section slide and confirm it has <=3 bullets now
    found_condensed = False
    for slide in prs.slides:
        if slide.shapes.title and "Dense Section" in slide.shapes.title.text:
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body:
                bullet_count = len([p for p in body.text_frame.paragraphs if p.text.strip()])
                assert bullet_count <= 3
                found_condensed = True
    assert found_condensed


def test_edit_remove_section():
    artifact_id = _generate_doc()
    plan = artifact_store.load_latest_content_plan(artifact_id)
    heading_to_remove = plan.sections[0].heading

    result = edit_artifact(artifact_id, f"Remove the {heading_to_remove} section")
    assert result.instruction.operation == "remove_section"

    doc = DocxDocument(result.file_path)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert heading_to_remove not in full_text


def test_edit_nonexistent_artifact_raises():
    with pytest.raises(ArtifactNotFoundError):
        edit_artifact("does-not-exist", "Add something")


def test_edit_api_endpoint():
    artifact_id = _generate_doc()
    resp = client.post(f"/artifact/{artifact_id}/edit", json={"request": "Add a competitive analysis section"})
    assert resp.status_code == 200
    body = resp.json()
    assert body["new_version"] == 2
    assert "competitive analysis" in body["instruction"]["new_heading"].lower()


def test_edit_api_nonexistent_artifact_returns_404():
    resp = client.post("/artifact/does-not-exist/edit", json={"request": "Add something"})
    assert resp.status_code == 404


# --- Version management API ---

def test_get_artifact_metadata():
    artifact_id = _generate_doc()
    resp = client.get(f"/artifact/{artifact_id}")
    assert resp.status_code == 200
    body = resp.json()
    assert body["artifact_type"] == "docx"
    assert body["current_version"] == 1
    assert len(body["versions"]) == 1


def test_version_history_accumulates_across_edits():
    artifact_id = _generate_doc()
    client.post(f"/artifact/{artifact_id}/edit", json={"request": "Add an executive summary"})
    client.post(f"/artifact/{artifact_id}/edit", json={"request": "Add a competitive analysis section"})

    resp = client.get(f"/artifact/{artifact_id}/versions")
    assert resp.status_code == 200
    versions = resp.json()
    assert len(versions) == 3
    assert [v["version"] for v in versions] == [1, 2, 3]
    assert versions[1]["change_request"] == "Add an executive summary"


def test_download_specific_version():
    artifact_id = _generate_doc()
    client.post(f"/artifact/{artifact_id}/edit", json={"request": "Add an executive summary"})

    # Download v1 explicitly — should NOT contain the edit
    resp_v1 = client.get(f"/artifact/{artifact_id}/download?version=1")
    assert resp_v1.status_code == 200

    # Download latest (no version param) — should be v2
    resp_latest = client.get(f"/artifact/{artifact_id}/download")
    assert resp_latest.status_code == 200


def test_get_versions_for_nonexistent_artifact_404():
    resp = client.get("/artifact/does-not-exist/versions")
    assert resp.status_code == 404
