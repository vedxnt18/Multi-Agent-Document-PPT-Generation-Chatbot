"""
backend/agents/ppt_analyzer.py

PPT Analysis Agent. Consumes an ExtractionResult for a PPTX and produces a
PPTStyleSpec. Slide dimensions/theme colors require reading the raw pptx
file directly (python-pptx object, not just our normalized blocks), so this
agent takes the original file path in addition to the extraction result.

content_density is a heuristic (blocks-per-slide thresholds), not an LLM
judgment — it's a measurable quantity. visual_style ("minimal corporate",
etc.) is the one LLM-classified field, same honesty principle as the
document analyzer: mock provider -> clearly labeled as unclassified.
"""
import logging
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.util import Emu

from backend.schemas.extraction import ExtractionResult, BlockType
from backend.schemas.template_spec import PPTStyleSpec, SlidePatternInfo, FontUsage
from backend.services.llm_service import get_llm_provider, LLMError, LLMJSONParseError

logger = logging.getLogger(__name__)

VISUAL_STYLE_SYSTEM_PROMPT = (
    "You are a presentation design assistant. Given slide titles and body "
    "text from a PowerPoint template, describe its visual/content style "
    "concisely."
)

DENSITY_LIGHT_MAX = 2
DENSITY_MODERATE_MAX = 5


def _extract_fonts_from_pptx(file_path: str) -> Counter:
    """Read actual run-level fonts directly from the pptx file (not available
    in our normalized blocks, which are text-only)."""
    counter: Counter = Counter()
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            counter[run.font.name] += 1
    except Exception as e:
        logger.warning(f"Could not extract font details from pptx: {e}")
    return counter


def _build_style_prompt(extraction: ExtractionResult) -> str:
    titles = [b.text for b in extraction.blocks if b.block_type == BlockType.SLIDE_TITLE][:12]
    bodies = [b.text for b in extraction.blocks if b.block_type == BlockType.SLIDE_BODY][:12]
    return (
        f"Slide titles: {titles}\n\n"
        f"Sample body content: {bodies}\n\n"
        'Return JSON with exactly this key: '
        '{"visual_style": "<short phrase, e.g. minimal corporate, bold and colorful, data-heavy technical>"}'
    )


def analyze_ppt(extraction: ExtractionResult, file_path: str, use_llm: bool = True) -> PPTStyleSpec:
    spec = PPTStyleSpec(file_id=extraction.file_id, slide_count=extraction.page_or_slide_count)

    if not extraction.blocks:
        spec.warnings.append("No content blocks found in extraction result; analysis is empty.")

    # --- Slide dimensions (from raw pptx, not in normalized blocks) ---
    try:
        prs = Presentation(file_path)
        spec.slide_width_in = float(Emu(prs.slide_width).inches)
        spec.slide_height_in = float(Emu(prs.slide_height).inches)
    except Exception as e:
        spec.warnings.append(f"Could not read slide dimensions: {e}")

    # --- Fonts (from raw pptx) ---
    font_counter = _extract_fonts_from_pptx(file_path)
    spec.fonts_used = [FontUsage(font_name=name, count=count) for name, count in font_counter.most_common()]
    spec.primary_font = spec.fonts_used[0].font_name if spec.fonts_used else None

    # --- Slide patterns + content density (structural, from normalized blocks) ---
    per_slide_bodies: dict[int, int] = defaultdict(int)
    per_slide_has_title: dict[int, bool] = defaultdict(bool)
    per_slide_has_table: dict[int, bool] = defaultdict(bool)
    per_slide_has_image: dict[int, bool] = defaultdict(bool)

    for block in extraction.blocks:
        sn = block.slide_number
        if sn is None:
            continue
        if block.block_type == BlockType.SLIDE_TITLE:
            per_slide_has_title[sn] = True
        elif block.block_type == BlockType.SLIDE_BODY:
            per_slide_bodies[sn] += 1
        elif block.block_type == BlockType.TABLE:
            per_slide_has_table[sn] = True
        elif block.block_type == BlockType.IMAGE:
            per_slide_has_image[sn] = True

    for sn in range(1, spec.slide_count + 1):
        spec.slide_patterns.append(
            SlidePatternInfo(
                slide_number=sn,
                has_title=per_slide_has_title.get(sn, False),
                body_block_count=per_slide_bodies.get(sn, 0),
                has_table=per_slide_has_table.get(sn, False),
                has_image=per_slide_has_image.get(sn, False),
            )
        )

    if spec.slide_count > 0:
        total_bodies = sum(per_slide_bodies.values())
        spec.avg_body_blocks_per_slide = round(total_bodies / spec.slide_count, 2)
        if spec.avg_body_blocks_per_slide <= DENSITY_LIGHT_MAX:
            spec.content_density = "light"
        elif spec.avg_body_blocks_per_slide <= DENSITY_MODERATE_MAX:
            spec.content_density = "moderate"
        else:
            spec.content_density = "dense"

    if extraction.warnings:
        spec.warnings.extend(extraction.warnings)

    # --- Subjective classification (LLM) ---
    if use_llm:
        try:
            provider = get_llm_provider()
            if provider.name == "mock":
                spec.visual_style = "unspecified (mock provider — no live classification)"
            else:
                prompt = _build_style_prompt(extraction)
                result = provider.generate_json(prompt, system_prompt=VISUAL_STYLE_SYSTEM_PROMPT)
                spec.visual_style = str(result.get("visual_style", "unspecified"))
        except (LLMError, LLMJSONParseError) as e:
            logger.warning(f"PPT visual style classification failed: {e}")
            spec.warnings.append(f"Visual style classification unavailable: {e}")

    return spec
