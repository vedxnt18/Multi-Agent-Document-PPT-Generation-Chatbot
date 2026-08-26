"""
tests/test_ppt_generation.py

Phase 9 tests. Verifies generate_pptx produces a REAL, openable .pptx with
correct structure — re-opened with python-pptx, not just checked to exist.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pytest
from pptx import Presentation
from fastapi.testclient import TestClient

from backend.main import app
from backend.agents.ppt_generator import generate_pptx, _section_to_bullets
from backend.schemas.content_plan import ContentPlan, GeneratedSection
from backend.schemas.template_spec import PPTStyleSpec
from backend.services.citation_service import citation_registry

client = TestClient(app)


@pytest.fixture(autouse=True)
def _reset_citations():
    citation_registry.reset()
    yield


def test_section_to_bullets_uses_existing_bullets():
    section = GeneratedSection(heading="X", bullet_points=["A", "B"])
    assert _section_to_bullets(section) == ["A", "B"]


def test_section_to_bullets_converts_paragraphs():
    section = GeneratedSection(heading="X", paragraphs=["First sentence. Second sentence.", "Another point here."])
    bullets = _section_to_bullets(section)
    assert len(bullets) == 2
    assert "First sentence" in bullets[0]


def test_generate_pptx_produces_real_openable_file(tmp_path):
    citation_registry.reset()
    cid = citation_registry.register("WEB", "Example Source", "https://example.com")

    plan = ContentPlan(
        title="Test Presentation",
        subtitle="Subtitle here",
        sections=[
            GeneratedSection(heading="Introduction", bullet_points=["Point A", "Point B"], citation_ids=[cid]),
            GeneratedSection(heading="Market Overview", bullet_points=["Trend 1", "Trend 2", "Trend 3"]),
        ],
    )

    output_path = str(tmp_path / "test.pptx")
    result_path = generate_pptx(plan, output_path)

    assert Path(result_path).exists()

    prs = Presentation(result_path)
    # Title slide + 2 content slides + sources slide = 4
    assert len(prs.slides) == 4

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    full_text = "\n".join(all_text)

    assert "Test Presentation" in full_text
    assert "Introduction" in full_text
    assert "Point A" in full_text
    assert "Market Overview" in full_text
    assert "Sources" in full_text
    assert cid in full_text


def test_generate_pptx_honors_slide_dimensions_from_style_spec(tmp_path):
    plan = ContentPlan(title="Dim Test", sections=[GeneratedSection(heading="S1", bullet_points=["b1"])])
    style_spec = PPTStyleSpec(file_id="x", slide_width_in=13.333, slide_height_in=7.5)

    output_path = str(tmp_path / "dim_test.pptx")
    generate_pptx(plan, output_path, style_spec=style_spec)

    prs = Presentation(output_path)
    from pptx.util import Inches
    assert abs(prs.slide_width - Inches(13.333)) < Inches(0.01)


def test_generate_pptx_pads_to_reach_target_slide_count(tmp_path):
    plan = ContentPlan(title="Short Deck", sections=[GeneratedSection(heading="Only Section", bullet_points=["one point"])])

    output_path = str(tmp_path / "padded.pptx")
    generate_pptx(plan, output_path, target_slide_count=6)

    prs = Presentation(output_path)
    assert len(prs.slides) == 6


def test_generate_pptx_splits_dense_section_across_slides(tmp_path):
    many_bullets = [f"Bullet {i}" for i in range(12)]
    plan = ContentPlan(title="Dense Deck", sections=[GeneratedSection(heading="Big Section", bullet_points=many_bullets)])

    output_path = str(tmp_path / "dense.pptx")
    generate_pptx(plan, output_path)

    prs = Presentation(output_path)
    # title + ceil(12/5)=3 content slides = 4 (no sources slide, no citations used)
    assert len(prs.slides) == 4


def test_generate_presentation_api_end_to_end():
    resp = client.post(
        "/generate/presentation",
        json={
            "request": "Create a presentation on AI trends",
            "research_findings": ["AI is growing fast [WEB-001]."],
            "slide_count": 5,
        },
    )
    assert resp.status_code == 200
    body = resp.json()
    assert body["filename"].endswith(".pptx")
    assert Path(body["file_path"]).exists()

    prs = Presentation(body["file_path"])
    assert len(prs.slides) == 5
